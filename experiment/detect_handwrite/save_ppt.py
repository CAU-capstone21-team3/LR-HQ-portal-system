import cv2
import io
import numpy as np
from PIL import Image
from pptx import Presentation # 라이브러리 
from pptx.util import Inches

# http://www.andrewjanowczyk.com/computationally-creating-a-powerpoint-presentation-of-experimental-results-using-python/
def addimagetoslide(slide,img,left,top, height, width, resize = 1.0):
    res = cv2.resize(img , None, fx=resize,fy=resize ,interpolation=cv2.INTER_CUBIC) #since the images are going to be small, we can resize them to prevent the final pptx file from being large for no reason
    image_stream = io.BytesIO()
    Image.fromarray(res).save(image_stream,format="PNG")
    slide.shapes.add_picture(image_stream, left, top ,height,width)
    image_stream.close()

def save_ppt(frame_array):

    prs = Presentation()
    
    layout = prs.slide_layouts[6]
    for frame in frame_array:
        slide = prs.slides.add_slide(layout)
        image = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        
        # 이미지 크기, 위치 일반화 추가할것 (현재 = 16:9기준)
        addimagetoslide(slide, image, Inches(0),Inches(0),Inches(10),Inches(5.625)) 
        

    prs.save('test.pptx')

